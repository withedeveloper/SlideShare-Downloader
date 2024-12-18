from flask import Flask, render_template, request, send_file, jsonify
import os
import requests
from bs4 import BeautifulSoup
import json
from PIL import Image
import zipfile
from pptx import Presentation
from pptx.util import Inches
import fitz  # PyMuPDF
import uuid
import shutil
import threading
import time

app = Flask(__name__)

# Function to clean up directories after download
def clean_up_directory(directory):
    attempts = 0
    while attempts < 5:  # Try for a maximum of 5 attempts
        try:
            shutil.rmtree(directory)
            break  # Successfully deleted, exit the loop
        except OSError as e:
            print(f"Error deleting directory {directory}: {e}")
            attempts += 1
            time.sleep(1)  # Wait for 1 second before retrying

    if attempts >= 5:
        print(f"Failed to delete directory {directory} after 5 attempts.")
    else:
        print(f"Successfully deleted directory {directory}.")

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/fetch_slides', methods=['POST'])
def fetch_slides():
    data = request.get_json()
    url = data['url']
    session_id = str(uuid.uuid4())

    try:
        response = requests.get(url)
        soup = BeautifulSoup(response.content, "html.parser")
        data = json.loads(soup.select_one("#__NEXT_DATA__").text)
        slides = data["props"]["pageProps"]["slideshow"]["slides"]
        total_slides = data["props"]["pageProps"]["slideshow"]["totalSlides"]

        host = slides["host"]
        image_location = slides["imageLocation"]
        image_sizes = slides["imageSizes"]
        title = slides["title"]
        quality = image_sizes[-1]["quality"]
        width = image_sizes[-1]["width"]

        thumbnails = []
        for i in range(1, total_slides + 1):
            img_url = f"{host}/{image_location}/{quality}/{title}-{i}-{width}.jpg"
            thumbnails.append(img_url)

        return jsonify({'slides': thumbnails, 'session_id': session_id})

    except Exception as e:
        print(f"Error fetching slides: {e}")
        return jsonify({'error': 'Failed to fetch slides'})

@app.route('/download_images', methods=['POST'])
def download_images():
    url = request.form.get('url')
    format_type = request.form.get('format')
    selected_slides = request.form.getlist('slides')
    session_id = request.form.get('session_id')

    def download_images_from_slideshare(url, selected_slides, session_id):
        try:
            response = requests.get(url)
            soup = BeautifulSoup(response.content, "html.parser")
            data = json.loads(soup.select_one("#__NEXT_DATA__").text)
            slides = data["props"]["pageProps"]["slideshow"]["slides"]
            total_slides = data["props"]["pageProps"]["slideshow"]["totalSlides"]

            host = slides["host"]
            image_location = slides["imageLocation"]
            image_sizes = slides["imageSizes"]
            title = slides["title"]
            quality = image_sizes[-1]["quality"]
            width = image_sizes[-1]["width"]

            save_directory = os.path.join('downloads', session_id, 'slides')
            os.makedirs(save_directory, exist_ok=True)

            selected_slides = [int(slide) for slide in selected_slides]

            for i in selected_slides:
                img_url = f"{host}/{image_location}/{quality}/{title}-{i}-{width}.jpg"
                filename = f"{i}.jpg"
                response = requests.get(img_url, stream=True)
                with open(os.path.join(save_directory, filename), 'wb') as f:
                    for chunk in response.iter_content(1024):
                        f.write(chunk)
                    # Update progress bar
                    total_downloaded = len([name for name in os.listdir(save_directory) if os.path.isfile(os.path.join(save_directory, name))])
                    with app.app_context():
                        app.config['PROGRESS'][session_id] = total_downloaded / len(selected_slides)

            return save_directory

        except Exception as e:
            print(f"Error downloading images: {e}")
            return None

    def create_pdf_from_images(directory):
        images = []
        for file in sorted(os.listdir(directory)):
            if file.endswith('.jpg'):
                image = Image.open(os.path.join(directory, file))
                images.append(image.convert('RGB'))
        if images:
            pdf_filename = f"{directory}.pdf"
            images[0].save(pdf_filename, save_all=True, append_images=images[1:])
            return pdf_filename
        return None

    def create_ppt_from_pdf(pdf_filename):
        doc = fitz.open(pdf_filename)
        ppt = Presentation()
        slide_width = ppt.slide_width
        slide_height = ppt.slide_height

        for page_number in range(len(doc)):
            page = doc.load_page(page_number)
            pix = page.get_pixmap()
            img_filename = f"slide_{page_number + 1}.png"
            pix.save(img_filename)

            slide_layout = ppt.slide_layouts[5]
            slide = ppt.slides.add_slide(slide_layout)

            slide.shapes.add_picture(img_filename, 0, 0, width=slide_width, height=slide_height)
            os.remove(img_filename)  # Cleanup the temporary image file

        ppt_filename = pdf_filename.replace('.pdf', '.pptx')
        ppt.save(ppt_filename)
        return ppt_filename

    save_directory = download_images_from_slideshare(url, selected_slides, session_id)

    if save_directory:
        if format_type == 'zip':
            zip_filename = f"{save_directory}.zip"
            with zipfile.ZipFile(zip_filename, 'w') as zipf:
                for file in sorted(os.listdir(save_directory)):
                    if file.endswith('.jpg'):
                        zipf.write(os.path.join(save_directory, file), file)
            cleanup_thread = threading.Thread(target=clean_up_directory, args=(os.path.dirname(save_directory),))
            cleanup_thread.start()
            return send_file(zip_filename, as_attachment=True)

        elif format_type == 'pdf':
            pdf_filename = create_pdf_from_images(save_directory)
            if pdf_filename:
                cleanup_thread = threading.Thread(target=clean_up_directory, args=(os.path.dirname(save_directory),))
                cleanup_thread.start()
                return send_file(pdf_filename, as_attachment=True)

        elif format_type == 'ppt':
            pdf_filename = create_pdf_from_images(save_directory)
            if pdf_filename:
                ppt_filename = create_ppt_from_pdf(pdf_filename)
                cleanup_thread = threading.Thread(target=clean_up_directory, args=(os.path.dirname(save_directory),))
                cleanup_thread.start()
                return send_file(ppt_filename, as_attachment=True)

        return "Invalid format specified."

    return "Failed to download images."

@app.route('/progress/<session_id>')
def progress(session_id):
    total_downloaded = app.config['PROGRESS'].get(session_id, 0)
    return jsonify({'progress': total_downloaded})

if __name__ == '__main__':
    app.config['PROGRESS'] = {}
    app.run(debug=True)
